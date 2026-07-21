"""把一个目录中的海报图片生成 16:9 自动循环播放 PPTX。"""

from __future__ import annotations

import argparse
import re
import sys
import tempfile
from pathlib import Path
from zipfile import ZipFile

from lxml import etree
from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from pptx import Presentation
from pptx.util import Inches


CANVAS_W, CANVAS_H = 1920, 1080
POSTER_MAX_W = int(CANVAS_W * 0.94)
POSTER_MAX_H = int(CANVAS_H * 0.93)
SUPPORTED_EXTENSIONS = {".jpg", ".jpeg", ".png"}
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"p": P_NS}


def natural_key(path: Path) -> list[object]:
    """让 2.jpg 排在 10.jpg 前，同时保持中文文件名可排序。"""
    return [int(part) if part.isdigit() else part.casefold() for part in re.split(r"(\d+)", path.name)]


def collect_posters(input_dir: Path) -> list[Path]:
    """只读取输入目录第一层，避免把输出或缓存再次当成素材。"""
    return sorted(
        (
            path
            for path in input_dir.iterdir()
            if path.is_file() and path.suffix.casefold() in SUPPORTED_EXTENSIONS
        ),
        key=natural_key,
    )


def open_rgb(path: Path) -> Image.Image:
    """读取图片、应用 EXIF 方向，并把透明区域安全铺成白色。"""
    with Image.open(path) as source:
        source.load()
        oriented = ImageOps.exif_transpose(source)
        if oriented.mode in {"RGBA", "LA"} or "transparency" in oriented.info:
            rgba = oriented.convert("RGBA")
            flattened = Image.new("RGBA", rgba.size, "white")
            flattened.alpha_composite(rgba)
            return flattened.convert("RGB")
        return oriented.convert("RGB")


def make_composite(poster_path: Path, output_path: Path) -> tuple[int, int]:
    """生成“完整海报居中 + 同源模糊暗背景”的 1920x1080 合成图。"""
    poster = open_rgb(poster_path)
    width, height = poster.size
    if width <= 0 or height <= 0:
        raise ValueError(f"图片尺寸无效：{poster_path.name}")

    background = ImageOps.fit(
        poster,
        (CANVAS_W, CANVAS_H),
        method=Image.Resampling.LANCZOS,
        centering=(0.5, 0.5),
    )
    background = background.filter(ImageFilter.GaussianBlur(radius=45))
    background = ImageEnhance.Brightness(background).enhance(0.52)
    background = ImageEnhance.Color(background).enhance(0.95)

    scale = min(POSTER_MAX_W / width, POSTER_MAX_H / height)
    display_size = (
        max(1, round(width * scale)),
        max(1, round(height * scale)),
    )
    foreground = poster.resize(display_size, Image.Resampling.LANCZOS)
    paste_at = (
        (CANVAS_W - display_size[0]) // 2,
        (CANVAS_H - display_size[1]) // 2,
    )
    background.paste(foreground, paste_at)
    background.save(output_path, "JPEG", quality=92, optimize=True, progressive=True)
    return width, height


def set_slide_transition(slide, seconds: float) -> None:
    """设置淡入淡出、按计时自动前进并禁用点击前进。"""
    slide_xml = slide._element
    transition = slide_xml.find("p:transition", namespaces=NSMAP)
    if transition is None:
        transition = etree.Element(f"{{{P_NS}}}transition")
        insert_before = slide_xml.find("p:timing", namespaces=NSMAP)
        if insert_before is None:
            insert_before = slide_xml.find("p:extLst", namespaces=NSMAP)
        if insert_before is None:
            slide_xml.append(transition)
        else:
            slide_xml.insert(slide_xml.index(insert_before), transition)

    transition.set("advTm", str(round(seconds * 1000)))
    transition.set("advClick", "0")
    for child in list(transition):
        transition.remove(child)
    etree.SubElement(transition, f"{{{P_NS}}}fade")


def set_loop_presentation(prs: Presentation) -> None:
    """把循环与使用计时写入正确的 Presentation Properties 部件。"""
    properties_part = next(
        (
            part
            for part in prs.part.package.iter_parts()
            if str(part.partname) == "/ppt/presProps.xml"
        ),
        None,
    )
    if properties_part is None:
        raise RuntimeError("PPTX 中缺少 ppt/presProps.xml")

    root = etree.fromstring(properties_part.blob)
    show_properties = root.find("p:showPr", namespaces=NSMAP)
    if show_properties is None:
        show_properties = etree.Element(f"{{{P_NS}}}showPr")
        extension_list = root.find("p:extLst", namespaces=NSMAP)
        if extension_list is None:
            root.append(show_properties)
        else:
            root.insert(root.index(extension_list), show_properties)

    show_properties.set("loop", "1")
    show_properties.set("useTimings", "1")
    for selector_name in ("custShow", "sldRg", "sldAll"):
        selector = show_properties.find(f"p:{selector_name}", namespaces=NSMAP)
        if selector is not None:
            show_properties.remove(selector)
    etree.SubElement(show_properties, f"{{{P_NS}}}sldAll")

    properties_part._blob = etree.tostring(
        root,
        encoding="UTF-8",
        xml_declaration=True,
        standalone=True,
    )


def verify_pptx(output_path: Path, expected_slides: int, seconds: float) -> None:
    """对关键 OOXML 设置做生成后自检，避免“生成成功但不会循环”。"""
    expected_ms = str(round(seconds * 1000))
    with ZipFile(output_path) as package:
        presentation = etree.fromstring(package.read("ppt/presentation.xml"))
        slide_size = presentation.find("p:sldSz", namespaces=NSMAP)
        if slide_size is None:
            raise RuntimeError("自检失败：缺少幻灯片尺寸")
        ratio = int(slide_size.get("cx")) / int(slide_size.get("cy"))
        if abs(ratio - (16 / 9)) > 0.0001:
            raise RuntimeError(f"自检失败：幻灯片不是 16:9（当前 {ratio:.4f}）")

        properties = etree.fromstring(package.read("ppt/presProps.xml"))
        show_properties = properties.find("p:showPr", namespaces=NSMAP)
        if show_properties is None or show_properties.get("loop") != "1":
            raise RuntimeError("自检失败：循环播放未写入 ppt/presProps.xml")
        if show_properties.get("useTimings") != "1":
            raise RuntimeError("自检失败：未启用幻灯片计时")

        slide_names = sorted(
            (
                name
                for name in package.namelist()
                if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
            ),
            key=lambda name: int(re.search(r"\d+", name).group()),
        )
        if len(slide_names) != expected_slides:
            raise RuntimeError(
                f"自检失败：预期 {expected_slides} 页，实际 {len(slide_names)} 页"
            )

        for slide_name in slide_names:
            slide = etree.fromstring(package.read(slide_name))
            transition = slide.find("p:transition", namespaces=NSMAP)
            if transition is None:
                raise RuntimeError(f"自检失败：{slide_name} 缺少切换设置")
            if transition.get("advTm") != expected_ms or transition.get("advClick") != "0":
                raise RuntimeError(f"自检失败：{slide_name} 自动切换设置不正确")
            if transition.find("p:fade", namespaces=NSMAP) is None:
                raise RuntimeError(f"自检失败：{slide_name} 缺少淡入淡出效果")


def build_presentation(
    poster_files: list[Path], output_path: Path, seconds: float
) -> list[tuple[Path, int, int]]:
    """生成 PPTX，并返回输入素材的尺寸记录。"""
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs._element.sldSz.set("type", "screen16x9")
    blank_layout = prs.slide_layouts[6]
    dimensions: list[tuple[Path, int, int]] = []

    with tempfile.TemporaryDirectory(prefix="campus-poster-display-") as temp_dir:
        temp_path = Path(temp_dir)
        for index, poster_path in enumerate(poster_files, 1):
            composite_path = temp_path / f"slide-{index}.jpg"
            try:
                width, height = make_composite(poster_path, composite_path)
            except Exception as exc:
                raise RuntimeError(f"处理图片失败：{poster_path.name}（{exc}）") from exc

            dimensions.append((poster_path, width, height))
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(composite_path),
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            set_slide_transition(slide, seconds)

        set_loop_presentation(prs)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)

    verify_pptx(output_path, len(poster_files), seconds)
    return dimensions


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="把目录中的 JPG/JPEG/PNG 海报生成 16:9 自动循环播放 PPTX。"
    )
    parser.add_argument(
        "--input",
        type=Path,
        default=Path.cwd(),
        help="海报目录，默认当前目录；只读取第一层图片",
    )
    parser.add_argument(
        "--output",
        type=Path,
        default=None,
        help="输出 PPTX，默认 <输入目录>/artifacts/poster-display.pptx",
    )
    parser.add_argument(
        "--seconds",
        type=float,
        default=10.0,
        help="每页停留秒数，默认 10",
    )
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    input_dir = args.input.expanduser().resolve()
    output_path = (
        args.output.expanduser().resolve()
        if args.output is not None
        else input_dir / "artifacts" / "poster-display.pptx"
    )

    if not input_dir.is_dir():
        print(f"ERROR: 输入目录不存在：{input_dir}", file=sys.stderr)
        return 2
    if args.seconds <= 0:
        print("ERROR: --seconds 必须大于 0", file=sys.stderr)
        return 2

    poster_files = collect_posters(input_dir)
    if not poster_files:
        print(
            f"ERROR: {input_dir} 第一层没有 JPG/JPEG/PNG 海报图片",
            file=sys.stderr,
        )
        return 1

    try:
        dimensions = build_presentation(poster_files, output_path, args.seconds)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 1

    print(f"已生成：{output_path}")
    print(f"幻灯片：{len(dimensions)} 页；每页：{args.seconds:g} 秒；循环：已启用")
    narrow_posters: list[str] = []
    for index, (path, width, height) in enumerate(dimensions, 1):
        print(f"  {index:02d}. {path.name} ({width}x{height})")
        if width / height < 0.55:
            narrow_posters.append(path.name)
    if narrow_posters:
        print(
            "WARN: 以下超长海报虽已完整保留，但远距离可读性可能不足；建议人工拆页："
            + "、".join(narrow_posters),
            file=sys.stderr,
        )
    print(f"文件大小：{output_path.stat().st_size / 1024:.0f} KB；OOXML 自检：通过")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
